# -*- coding: utf-8 -*-
"""
@Time ： 2021/12/31 0:37
@Auth ： 张顺
@No   : 021321712238
@File ：pptReader.py
@IDE ：PyCharm

"""
import os
import re

from pptx import Presentation

def getFileOrDirPath(name):
    return os.path.join(os.path.dirname(__file__), name)

if __name__ == "__main__":
    mine = ''
    for i in os.listdir(getFileOrDirPath("ppts")):
        FILE_PATH = getFileOrDirPath("ppts/"+ i)
        pptx = Presentation(FILE_PATH)
        for slide in pptx.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        mine += re.sub("\d{2}:\d{2}:\d{2}","",paragraph.text.replace("\n",""))
            mine += "\n"
                        # for run in paragraph.runs:
                        #     print(run.text)
        print(mine)
    with open(getFileOrDirPath("txt/计算机网络第8版课件.txt"), "w+",encoding="utf-8") as f:
        f.writelines(mine)